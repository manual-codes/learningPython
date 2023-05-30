import os
import docx
from openpyxl import Workbook
from pptx import Presentation

def create_files():
    # Get the current directory
    current_dir = os.path.dirname(os.path.realpath(__file__))

    # Get the name of the directory
    directory_name = os.path.basename(current_dir)

    # Prefix for the files
    prefix = directory_name

    # File details (type, suffix)
    file_details = [
        ("Word", "SOP"),
        ("Word", "ODI_PDI"),
        ("Excel", "PIV Sheet"),
        ("Excel", "Production Errors"),
        ("Excel", "Project Plan"),
        ("Excel", "CBA"),
        ("PowerPoint", "Process Map"),
        ("Excel", "Ramp Up Volume Plan")
    ]

    for file_type, suffix in file_details:
        # Construct the file name
        file_name = f"{prefix}_{suffix}"

        # Create the file based on the type
        if file_type == "Word":
            create_word_file(current_dir, file_name)
        elif file_type == "Excel":
            create_excel_file(current_dir, file_name)
        elif file_type == "PowerPoint":
            create_powerpoint_file(current_dir, file_name)

    print("Files created successfully!")

def create_word_file(directory, file_name):
    # Create a Word document
    doc = docx.Document()
    doc.add_paragraph(f"This is a sample {file_name} Word document.")

    # Save the Word document
    doc.save(os.path.join(directory, file_name + ".docx"))

def create_excel_file(directory, file_name):
    # Create an Excel file
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = file_name
    sheet["A1"] = f"Sample {file_name}"

    # Save the Excel file
    workbook.save(os.path.join(directory, file_name + ".xlsx"))

def create_powerpoint_file(directory, file_name):
    # Create a PowerPoint presentation
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[0])
    title = slide.shapes.title
    title.text = f"Sample {file_name}"

    # Save the PowerPoint presentation
    presentation.save(os.path.join(directory, file_name + ".pptx"))

if __name__ == "__main__":
    create_files()
